#!/usr/bin/env python3
"""
Generate a 10-minute thesis defense presentation for TCN-HMAC.
Outputs: thesis_defense_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────
DARK_BG      = RGBColor(0x1B, 0x1B, 0x2F)   # dark navy background
ACCENT_BLUE  = RGBColor(0x00, 0x7A, 0xCC)   # accent blue
ACCENT_GREEN = RGBColor(0x00, 0xB3, 0x6B)   # accent green for success metrics
ACCENT_RED   = RGBColor(0xE0, 0x4F, 0x5F)   # accent red for problems
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
GOLD         = RGBColor(0xFF, 0xD7, 0x00)
SOFT_BG      = RGBColor(0xF5, 0xF7, 0xFA)   # light slide bg
HEADER_BG    = RGBColor(0x00, 0x3F, 0x72)   # dark blue header bar
TEXT_DARK    = RGBColor(0x22, 0x22, 0x22)
TEXT_MED     = RGBColor(0x44, 0x44, 0x44)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FIGURES = os.path.join(os.path.dirname(__file__), "Figures")
RESULTS = os.path.join(FIGURES, "results")

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── Helper functions ────────────────────────────────────────
def add_bg(slide, color=SOFT_BG):
    """Fill entire slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title_text, subtitle_text=None):
    """Add a coloured header bar across the top."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = HEADER_BG
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # left padding via space indent
    p.space_before = Pt(6)

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = LIGHT_GRAY
        p2.alignment = PP_ALIGN.LEFT

    # Slide number placeholder (bottom right)
    num_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(1.2), SLIDE_HEIGHT - Inches(0.45),
        Inches(1), Inches(0.35)
    )
    num_tf = num_box.text_frame
    num_tf.paragraphs[0].text = ""  # will be filled later
    return bar


def add_bullet_textbox(slide, left, top, width, height, items,
                       font_size=18, color=TEXT_DARK, bold_prefix=True,
                       line_spacing=1.4, bullet_char="\u25B8"):
    """Add a textbox with bulleted items. Items can be str or (bold_part, rest)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = Pt(int(font_size * line_spacing))

        if isinstance(item, tuple):
            # (bold_text, normal_text)
            run_b = p.add_run()
            run_b.text = f"{bullet_char} {item[0]}"
            run_b.font.size = Pt(font_size)
            run_b.font.bold = True
            run_b.font.color.rgb = color

            run_n = p.add_run()
            run_n.text = f" {item[1]}"
            run_n.font.size = Pt(font_size)
            run_n.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = f"{bullet_char} {item}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return txBox


def add_metric_box(slide, left, top, width, height,
                   value_text, label_text, value_color=ACCENT_GREEN):
    """Add a rounded-rectangle metric card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value_text
    r1.font.size = Pt(30)
    r1.font.bold = True
    r1.font.color.rgb = value_color

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label_text
    r2.font.size = Pt(13)
    r2.font.color.rgb = TEXT_MED
    return shape


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add image if it exists, else add placeholder text."""
    if os.path.isfile(img_path):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(img_path, **kwargs)
    else:
        txBox = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.5))
        txBox.text_frame.paragraphs[0].text = f"[Image not found: {os.path.basename(img_path)}]"


def set_slide_numbers():
    """Write slide numbers into the last textbox of each slide."""
    for idx, slide in enumerate(prs.slides):
        for shape in reversed(slide.shapes):
            if shape.has_text_frame and shape.text == "":
                shape.text_frame.paragraphs[0].text = str(idx + 1)
                shape.text_frame.paragraphs[0].font.size = Pt(11)
                shape.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                break


# ════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# Title
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(2.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "TCN-HMAC"
r.font.size = Pt(54)
r.font.bold = True
r.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "A Lightweight Deep Learning and Cryptographic\nHybrid Security Framework for SDN"
r2.font.size = Pt(26)
r2.font.color.rgb = LIGHT_GRAY

# Subtitle line
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.9), Inches(3.3), Pt(3)
)
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Author info
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.3), Inches(2))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p3 = tf2.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "Md. Mahamudul Hasan Zubayer"
r3.font.size = Pt(22)
r3.font.color.rgb = WHITE

p4 = tf2.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
r4 = p4.add_run()
r4.text = "ID: 2022000000006"
r4.font.size = Pt(16)
r4.font.color.rgb = LIGHT_GRAY

p5 = tf2.add_paragraph()
p5.alignment = PP_ALIGN.CENTER
p5.space_before = Pt(16)
r5 = p5.add_run()
r5.text = "Supervisor: Dr. Md. Maruf Hassan, Associate Professor"
r5.font.size = Pt(16)
r5.font.color.rgb = LIGHT_GRAY

p6 = tf2.add_paragraph()
p6.alignment = PP_ALIGN.CENTER
p6.space_before = Pt(4)
r6 = p6.add_run()
r6.text = "Department Head: Shahriar Manzoor, Associate Professor"
r6.font.size = Pt(16)
r6.font.color.rgb = LIGHT_GRAY

p7 = tf2.add_paragraph()
p7.alignment = PP_ALIGN.CENTER
p7.space_before = Pt(16)
r7 = p7.add_run()
r7.text = "Department of Computer Science and Engineering\nSoutheast University"
r7.font.size = Pt(16)
r7.font.color.rgb = LIGHT_GRAY

# University logo
logo_path = os.path.join(FIGURES, "SEULogo.png")
if os.path.isfile(logo_path):
    slide.shapes.add_picture(logo_path, Inches(6.0), Inches(6.2), height=Inches(0.9))

# ════════════════════════════════════════════════════════════
# SLIDE 2 — Outline
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Presentation Outline")

outline_items = [
    "Background & Motivation",
    "Problem Statement",
    "Research Objectives",
    "Proposed Methodology — TCN-HMAC Framework",
    "Dataset & Preprocessing",
    "TCN Model Architecture",
    "Experimental Results & Analysis",
    "Comparative Analysis",
    "Conclusion & Future Work",
]
add_bullet_textbox(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(5.5),
                   outline_items, font_size=22, line_spacing=1.55,
                   bullet_char="\u25CF")


# ════════════════════════════════════════════════════════════
# SLIDE 3 — Background & Motivation
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Background & Motivation")

items = [
    ("SDN Architecture:", "Separates control plane from data plane; a single controller orchestrates all switches via OpenFlow."),
    ("Centralization Risk:", "The controller becomes a single point of failure\u2014its compromise cascades across the entire network."),
    ("Control Channel Vulnerability:", "OpenFlow messages (flow rules, stats) can be intercepted, modified, or injected if not properly secured."),
    ("Flow Rule Integrity Gap:", "Switches blindly execute any received flow rule\u2014no built-in mechanism to verify authenticity."),
    ("Growing Attack Sophistication:", "Modern adversaries chain DDoS, MITM, probe, and brute-force attacks simultaneously."),
    ("TLS Limitations:", "TLS protects transport only; it cannot verify flow rule semantics or stop a compromised controller."),
]
add_bullet_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                   items, font_size=18, line_spacing=1.5)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — Problem Statement
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Problem Statement")

problems = [
    ("Detection without Verification:", "ML/DL-based IDS can flag malicious traffic but cannot prevent tampered flow rules from being installed."),
    ("Verification without Detection:", "Cryptographic mechanisms (HMAC, TLS) verify rule integrity but do not proactively detect intrusion attempts."),
    ("Excessive Computational Overhead:", "Blockchain/PKI-based solutions introduce latency incompatible with real-time SDN operations."),
    ("Limited Attack Coverage:", "Most solutions target a single attack vector (e.g., DDoS only), lacking multi-threat generality."),
    ("Deployment Complexity:", "Solutions requiring OpenFlow protocol changes or custom switch firmware hinder real-world adoption."),
]
add_bullet_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.0),
                   problems, font_size=18, color=TEXT_DARK, line_spacing=1.5)

# Research question box
rq_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.2), Inches(5.6), Inches(10.9), Inches(1.4)
)
rq_box.fill.solid()
rq_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
rq_box.line.color.rgb = ACCENT_BLUE
rq_box.line.width = Pt(2)

tf = rq_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Research Question: "
r.font.size = Pt(15)
r.font.bold = True
r.font.color.rgb = HEADER_BG
rn = p.add_run()
rn.text = "How can we build a lightweight, deployable security framework for SDN that combines proactive multi-class intrusion detection with real-time cryptographic verification\u2014without exceeding the latency budget that production networks demand?"
rn.font.size = Pt(15)
rn.font.color.rgb = TEXT_DARK
rn.font.italic = True


# ════════════════════════════════════════════════════════════
# SLIDE 5 — Research Objectives
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Research Objectives")

objectives = [
    ("RO1:", "Design a TCN-based IDS achieving >99% accuracy for binary classification of DDoS, MITM, Probe, and Brute-force attacks."),
    ("RO2:", "Develop a preprocessing pipeline (Pearson correlation, PCA, class weighting) for the InSDN dataset."),
    ("RO3:", "Train & deploy the TCN model via TensorFlow/Keras with compact size suitable for resource-constrained environments."),
    ("RO4:", "Design a lightweight auxiliary agent for HMAC-based flow rule verification and challenge\u2013response controller authentication."),
    ("RO5:", "Integrate the TCN-IDS and HMAC agent into a cohesive framework within the SDN control loop."),
    ("RO6:", "Evaluate the framework on accuracy, precision, recall, F1-score, and computational overhead."),
    ("RO7:", "Conduct comparative analysis against 15 existing models (LSTM, CNN-BiLSTM, Transformer, etc.)."),
]
add_bullet_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                   objectives, font_size=17, line_spacing=1.45)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — Proposed Methodology / System Architecture
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Proposed Methodology: TCN-HMAC Framework")

# Left column - key points
items = [
    ("Layer 1 \u2014 TCN-IDS:", "Deployed as a controller app; inspects flow statistics in real time; classifies traffic as benign or malicious."),
    ("Layer 2 \u2014 HMAC Agent:", "Co-located subprocess; HMAC-SHA256 verification of every flow rule; challenge\u2013response for controller authenticity."),
    ("Defense in Depth:", "If an attack bypasses the TCN, the HMAC prevents rogue flow rules. If keys are compromised, the TCN detects malicious patterns."),
]
add_bullet_textbox(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(3.5),
                   items, font_size=16, line_spacing=1.45)

# System architecture image
arch_img = os.path.join(FIGURES, "system_architecture.drawio.png")
add_image_safe(slide, arch_img, Inches(6.5), Inches(1.4), width=Inches(6.3))

# Bottom note
note = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1.2))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Key: No OpenFlow protocol modifications required. Runs entirely in software as a standard controller application + auxiliary agent."
r.font.size = Pt(14)
r.font.italic = True
r.font.color.rgb = TEXT_MED


# ════════════════════════════════════════════════════════════
# SLIDE 7 — Dataset & Preprocessing
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Dataset & Preprocessing Pipeline")

items = [
    ("InSDN Dataset:", "Purpose-built for SDN IDS research; covers Normal, DDoS, MITM, Probe, Brute-force traffic."),
    ("Raw:", "343,889 samples \u00d7 84 features \u2192 Cleaned: 182,831 samples \u00d7 48 features."),
    ("Feature Selection:", "Pearson correlation (|r| > 0.95 removal) \u2192 removed redundant highly-correlated features."),
    ("Normalization:", "StandardScaler (zero mean, unit variance) for stable gradient-based training."),
    ("Dimensionality Reduction:", "PCA: 48 \u2192 24 features, retaining 95.43% of total variance."),
    ("Class Weighting:", "Inverse-frequency weighting to counter class imbalance (65% attack, 35% benign)."),
    ("Train/Test Split:", "80/20 stratified split \u2192 146,264 training / 36,567 test samples."),
]
add_bullet_textbox(slide, Inches(0.5), Inches(1.4), Inches(7.0), Inches(5.5),
                   items, font_size=16, line_spacing=1.45)

# PCA variance plot
pca_img = os.path.join(RESULTS, "pca_explained_variance.png")
add_image_safe(slide, pca_img, Inches(7.8), Inches(1.6), width=Inches(5.2))

# Class distribution plot
cd_img = os.path.join(RESULTS, "eda_class_distribution.png")
add_image_safe(slide, cd_img, Inches(7.8), Inches(4.5), width=Inches(5.2))


# ════════════════════════════════════════════════════════════
# SLIDE 8 — TCN Model Architecture
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "TCN Model Architecture")

arch_items = [
    ("Input:", "24-dimensional PCA-transformed feature vector reshaped to (24, 1)."),
    ("6 Residual Blocks:", "Dilated causal Conv1D layers with dilation rates [1, 2, 4, 8, 16, 32], 64 filters, kernel size 3."),
    ("Regularization:", "Batch Normalization + Spatial Dropout (20%) in each block."),
    ("Pooling:", "Global Average Pooling to aggregate temporal features into a fixed-length vector."),
    ("Classification Head:", "Dense(64, ReLU) \u2192 Dropout(30%) \u2192 Dense(1, Sigmoid) for binary output."),
    ("Total Parameters:", "156,737 (612 KB model size)\u2014compact enough for edge deployment."),
    ("Training:", "Adam optimizer, lr=10\u207b\u00b3 with ReduceLROnPlateau; binary cross-entropy; ~30 epochs (~5 min on T4 GPU)."),
]
add_bullet_textbox(slide, Inches(0.5), Inches(1.4), Inches(12), Inches(5.5),
                   arch_items, font_size=17, line_spacing=1.45)

# Why TCN box
why_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2)
)
why_box.fill.solid()
why_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF0)
why_box.line.color.rgb = ACCENT_GREEN
why_box.line.width = Pt(1.5)
tf = why_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Why TCN? "
r.font.size = Pt(15)
r.font.bold = True
r.font.color.rgb = ACCENT_GREEN
r2 = p.add_run()
r2.text = "Fully parallelizable inference (unlike sequential LSTM/GRU) \u2022 Stable gradients via residual connections \u2022 Controllable receptive field via dilation \u2022 5\u201330\u00d7 smaller than comparable DL models"
r2.font.size = Pt(14)
r2.font.color.rgb = TEXT_DARK


# ════════════════════════════════════════════════════════════
# SLIDE 9 — Results: Key Metrics (the hero slide)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Experimental Results", "Test Set: 36,567 samples (20% held-out, never seen during training)")

# Metric cards - top row
metrics = [
    ("99.85%", "Accuracy"),
    ("99.80%", "Precision"),
    ("99.97%", "Recall / DR"),
    ("99.89%", "F1-Score"),
    ("0.9999", "AUC-ROC"),
    ("0.37%", "False Alarm Rate"),
]
card_w = Inches(1.85)
card_h = Inches(1.2)
start_x = Inches(0.5)
gap = Inches(0.18)
y_pos = Inches(1.5)

for i, (val, label) in enumerate(metrics):
    clr = ACCENT_GREEN if i < 5 else ACCENT_RED
    add_metric_box(slide, start_x + i * (card_w + gap), y_pos,
                   card_w, card_h, val, label, value_color=clr)

# Confusion matrix image
cm_img = os.path.join(RESULTS, "plot_confusion_matrix.png")
add_image_safe(slide, cm_img, Inches(0.5), Inches(3.1), height=Inches(3.8))

# Right side - key takeaways
takeaways = [
    ("23,737 / 23,744", "attacks correctly detected (only 7 missed)."),
    ("47 / 12,823", "benign flows false-alarmed (0.37% rate)."),
    ("Error Bias:", "87% of errors are FP (conservative); only 13% are FN."),
    ("Operational:", "~1 false alarm per 270 benign flows\u2014manageable with whitelisting."),
]
add_bullet_textbox(slide, Inches(5.5), Inches(3.3), Inches(7.5), Inches(4.0),
                   takeaways, font_size=17, line_spacing=1.6)


# ════════════════════════════════════════════════════════════
# SLIDE 10 — Results: Training Curves & ROC
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Training Convergence & ROC Analysis")

# Accuracy plot
acc_img = os.path.join(RESULTS, "plot_accuracy.png")
add_image_safe(slide, acc_img, Inches(0.3), Inches(1.4), width=Inches(4.1))

# Loss plot
loss_img = os.path.join(RESULTS, "plot_loss.png")
add_image_safe(slide, loss_img, Inches(4.6), Inches(1.4), width=Inches(4.1))

# ROC curve
roc_img = os.path.join(RESULTS, "plot_roc_curve.png")
add_image_safe(slide, roc_img, Inches(8.9), Inches(1.4), width=Inches(4.1))

# Annotations
notes = [
    "Converges in ~30 epochs (5 min on T4 GPU)",
    "Train/Val curves align closely \u2192 no overfitting",
    "AUC-ROC = 0.9999 \u2192 near-perfect discrimination",
]
add_bullet_textbox(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.8),
                   notes, font_size=16, line_spacing=1.5, bullet_char="\u2713",
                   color=ACCENT_GREEN)


# ════════════════════════════════════════════════════════════
# SLIDE 11 — HMAC Overhead Results
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "HMAC Verification Overhead")

hmac_metrics = [
    ("~2 \u00b5s", "Per-Message\nHMAC Compute"),
    ("0.7 ms", "Total Control\nLatency Added"),
    ("4.4%", "CPU Usage\nOverhead"),
    ("13.7 MB", "RAM\nOverhead"),
    ("32 B", "Per-Message\nSize Overhead"),
]
card_w = Inches(2.1)
card_h = Inches(1.4)
start_x = Inches(0.7)
gap = Inches(0.22)
y_pos = Inches(1.6)

for i, (val, label) in enumerate(hmac_metrics):
    add_metric_box(slide, start_x + i * (card_w + gap), y_pos,
                   card_w, card_h, val, label, value_color=ACCENT_BLUE)

hmac_features = [
    ("Message Integrity:", "HMAC-SHA256 tag appended to every Flow_Mod message; switch-side verification via shadow table."),
    ("Replay Prevention:", "Sequence numbers + timestamps reject replayed or out-of-order messages."),
    ("Controller Authentication:", "Periodic challenge\u2013response protocol confirms the controller\u2019s identity."),
    ("Key Management:", "Pre-shared symmetric keys with planned rotation; no PKI infrastructure required."),
    ("Negligible Impact:", "Combined TCN inference + HMAC verification < 1 ms per flow\u2014well within real-time budgets."),
]
add_bullet_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.5),
                   hmac_features, font_size=17, line_spacing=1.5)


# ════════════════════════════════════════════════════════════
# SLIDE 12 — Comparative Analysis
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Comparative Analysis", "Against 15 existing IDS models & SDN security frameworks")

# Comparison table
from pptx.util import Emu as E

table_data = [
    ["Model", "Dataset", "Acc.(%)", "Prec.(%)", "Recall(%)", "F1(%)", "Params", "Inf. Time"],
    ["TCN-HMAC (Ours)", "InSDN", "99.85", "99.80", "99.97", "99.89", "157K", "0.17ms"],
    ["CNN-BiLSTM", "InSDN", "99.90", "99.91", "99.90", "99.90", "~2M", "~2ms"],
    ["DNN Ensemble", "InSDN", "99.70", "99.65", "99.70", "99.67", "~5M", "~3ms"],
    ["CNN-LSTM", "CIC-IDS", "99.67", "99.68", "99.67", "99.67", "~1.5M", "~1.5ms"],
    ["DRL (DDQN)", "InSDN", "98.85", "98.90", "98.85", "98.87", "~3M", "~5ms"],
    ["BiTCN-MHSA", "CIC-IDS", "99.72", "99.68", "99.72", "99.70", "~500K", "~0.5ms"],
    ["LSTM", "NSL-KDD", "99.23", "99.00", "99.23", "99.11", "~500K", "~1ms"],
    ["DT/RF Baseline", "InSDN", "98.50", "98.40", "98.50", "98.45", "\u2014", "\u2014"],
]

rows = len(table_data)
cols = len(table_data[0])
tbl_shape = slide.shapes.add_table(rows, cols,
                                    Inches(0.3), Inches(1.4),
                                    Inches(12.7), Inches(4.2))
tbl = tbl_shape.table

# Style the table
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(13)
                if r == 0:
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                elif r == 1:
                    run.font.bold = True
                    run.font.color.rgb = HEADER_BG
                else:
                    run.font.color.rgb = TEXT_DARK

        # Header row bg
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_BG
        elif r == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE0, 0xF0, 0xFF)
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# Key advantages at bottom
adv_items = [
    ("Highest Detection Rate:", "99.97% (70% reduction in miss rate vs next-best)."),
    ("Smallest Model:", "612 KB / 157K params\u20145\u201330\u00d7 smaller than comparable DL models."),
    ("Fastest Inference:", "0.17 ms\u2014enables 5,000+ flow classifications/sec."),
    ("Only Dual-Layer Framework:", "Combines DL-based IDS + cryptographic flow rule authentication."),
]
add_bullet_textbox(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(1.5),
                   adv_items, font_size=14, line_spacing=1.35, color=ACCENT_GREEN)


# ════════════════════════════════════════════════════════════
# SLIDE 13 — Research Objectives Fulfilled
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Objectives Fulfilled", "Mapping results back to research objectives")

fulfilled = [
    ("RO1 \u2714", "TCN achieves 99.85% accuracy, 99.97% recall across DDoS, MITM, Probe, Brute-force (target: >99%)."),
    ("RO2 \u2714", "15-stage pipeline: cleaning \u2192 Pearson correlation \u2192 StandardScaler \u2192 PCA (48\u219224 features, 95.43% variance)."),
    ("RO3 \u2714", "TensorFlow/Keras model: 612 KB, 156,737 parameters, 0.17 ms inference\u2014portable across platforms."),
    ("RO4 \u2714", "Auxiliary agent: HMAC-SHA256 verification (~2 \u00b5s/msg) + challenge\u2013response authentication (0.7 ms overhead)."),
    ("RO5 \u2714", "Integrated TCN-IDS + HMAC agent as Ryu controller app + co-located subprocess."),
    ("RO6 \u2714", "Full evaluation: Accuracy 99.85%, Precision 99.80%, Recall 99.97%, F1 99.89%, FAR 0.37%."),
    ("RO7 \u2714", "Benchmarked against 15 models\u2014highest detection rate, smallest model, fastest inference."),
]
add_bullet_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
                   fulfilled, font_size=16, line_spacing=1.42, bullet_char="")


# ════════════════════════════════════════════════════════════
# SLIDE 14 — Contributions
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Key Contributions")

contribs = [
    ("Novel TCN-HMAC Framework:", "First to combine temporal DL with lightweight cryptographic verification in a unified, deployable SDN security solution."),
    ("First TCN on InSDN:", "Highest detection rate (99.97%) and lowest FAR (0.37%) among 15 compared approaches."),
    ("Compact Deployment:", "612 KB model (157K params)\u2014runs inside the SDN controller without additional hardware."),
    ("Lightweight HMAC Agent:", "~2 \u00b5s per message, 0.7 ms control latency, 4.4% CPU, 13.7 MB RAM\u2014negligible overhead."),
    ("Comprehensive Benchmark:", "Most extensive comparative analysis in SDN IDS literature (15 models, multiple metrics)."),
]
add_bullet_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                   contribs, font_size=18, line_spacing=1.55)


# ════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusion & Future Work
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Conclusion & Future Work")

# Conclusion points
conc = [
    "TCN-HMAC delivers near-perfect IDS (99.85% acc, 99.97% DR, 0.37% FAR) in a 612 KB, sub-millisecond package.",
    "Bridges the detection\u2013verification divide: TCN handles data-plane threats; HMAC secures the control plane.",
    "Deployable as a standard controller app\u2014no OpenFlow changes, no custom hardware.",
    "Outperforms or matches 15 state-of-the-art models while being 5\u201330\u00d7 smaller and fastest at inference.",
]
conc_box = add_bullet_textbox(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.5),
                               conc, font_size=16, line_spacing=1.5, bullet_char="\u2713",
                               color=ACCENT_GREEN)

# Future work
future = [
    "Multi-dataset evaluation (NSL-KDD, CIC-IDS-2017, UNSW-NB15)",
    "Multi-class classification (identify specific attack types)",
    "Adversarial robustness evaluation (FGSM, PGD)",
    "Federated learning for cross-domain privacy-preserving training",
    "Online / incremental learning for concept drift adaptation",
    "Production HMAC implementation & hardware benchmarking",
    "Model compression (INT8 quantization, pruning) for edge deployment",
]

# Future work heading
fw_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(0.5))
tf = fw_title.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Future Work Directions"
r.font.size = Pt(20)
r.font.bold = True
r.font.color.rgb = HEADER_BG

add_bullet_textbox(slide, Inches(6.8), Inches(2.0), Inches(6), Inches(4.5),
                   future, font_size=15, line_spacing=1.45, bullet_char="\u25CB",
                   color=TEXT_MED)

# Bottom bar
bottom = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.7), SLIDE_WIDTH, Inches(0.8)
)
bottom.fill.solid()
bottom.fill.fore_color.rgb = HEADER_BG
bottom.line.fill.background()

tf = bottom.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "\"TCN-HMAC narrows the gap between what the SDN control plane needs and what existing solutions provide.\""
r.font.size = Pt(16)
r.font.italic = True
r.font.color.rgb = WHITE


# ════════════════════════════════════════════════════════════
# SLIDE 16 — Thank You / Q&A
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
r1 = p1.add_run()
r1.text = "Thank You"
r1.font.size = Pt(54)
r1.font.bold = True
r1.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)
r2 = p2.add_run()
r2.text = "Questions & Discussion"
r2.font.size = Pt(28)
r2.font.color.rgb = LIGHT_GRAY

# Decorative line
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5), Inches(5.0), Inches(3.3), Pt(3)
)
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Contact info
txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.4), Inches(11.3), Inches(1.2))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p = tf3.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Md. Mahamudul Hasan Zubayer  \u2022  Supervisor: Dr. Md. Maruf Hassan"
r.font.size = Pt(16)
r.font.color.rgb = LIGHT_GRAY

p2 = tf3.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Department of Computer Science and Engineering, Southeast University"
r2.font.size = Pt(14)
r2.font.color.rgb = LIGHT_GRAY


# ── Finalize ────────────────────────────────────────────────
set_slide_numbers()
out_path = os.path.join(os.path.dirname(__file__), "thesis_defense_presentation.pptx")
prs.save(out_path)
print(f"Presentation saved to: {out_path}")
print(f"Total slides: {len(prs.slides)}")
